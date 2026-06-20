"""PDF, DOCX, and PPTX loader.

Purpose
-------
Reads office and PDF documents and extracts their text, tables, and (where needed)
image content into normalized records.

What it does
------------
Loads PDFs page by page, extracting native text where present and falling back to
OCR or vision captioning for image-only pages; loads DOCX text and tables; loads
PPTX slide text.

Flow
----
For each page or section it extracts text directly when available, otherwise tries
OCR, then optional vision captioning, and as a last resort catalogs the page so it
is still represented. Tables are flattened into readable rows.
"""

from pathlib import Path
from typing import Any

import fitz
from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from backend.app.core.config import VisionConfig
from backend.app.ingestion.loaders.base import LoadedFile, LoadedRecord
from backend.app.vision.local_vision_client import LocalVisionClient
from backend.app.vision.caption_budget import CaptionBudget
from backend.app.core.config import OcrConfig, VisionConfig
from backend.app.ocr.local_ocr import LocalOcrClient
from backend.app.ocr.ocr_budget import OcrBudget


class DocumentFileLoader:
    """
    Loader for PDF, DOCX, and PPTX files.
    """

    supported_extensions = {
        ".pdf",
        ".docx",
        ".pptx",
    }

    pdf_extensions = {".pdf"}
    docx_extensions = {".docx"}
    pptx_extensions = {".pptx"}

    def __init__(
        self,
        vision_config: VisionConfig | None = None,
        caption_budget: CaptionBudget | None = None,
        ocr_config: OcrConfig | None = None,
        ocr_budget: OcrBudget | None = None,
    ):
        self.vision_config = vision_config or VisionConfig()
        self.caption_budget = caption_budget or CaptionBudget.from_config(self.vision_config)
        self.vision_client = LocalVisionClient(self.vision_config)

        self.ocr_config = ocr_config or OcrConfig()
        self.ocr_budget = ocr_budget or OcrBudget.from_config(self.ocr_config)
        self.ocr_client = LocalOcrClient(self.ocr_config)

    def supports(self, path: str | Path) -> bool:
        return Path(path).suffix.lower() in self.supported_extensions

    def load(self, path: str | Path) -> LoadedFile:
        file_path = Path(path)

        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        suffix = file_path.suffix.lower()

        if suffix in self.pdf_extensions:
            return self._load_pdf(file_path)

        if suffix in self.docx_extensions:
            return self._load_docx(file_path)

        if suffix in self.pptx_extensions:
            return self._load_pptx(file_path)

        raise ValueError(f"Unsupported document file type: {suffix}")

    def _load_pdf(self, path: Path) -> LoadedFile:
        """
        Load a PDF.

        For each page:
            - use embedded text if available
            - otherwise use optional vision captioning
            - otherwise create image-page metadata
        """

        reader = PdfReader(str(path))
        records: list[LoadedRecord] = []

        for page_index, page in enumerate(reader.pages):
            page_number = page_index + 1
            text = page.extract_text() or ""
            text = text.strip()

            if text:
                records.append(
                    LoadedRecord(
                        text=text,
                        record_type="pdf_page",
                        title=f"{path.name} page {page_number}",
                        metadata={
                            "file_name": path.name,
                            "file_suffix": path.suffix.lower(),
                            "page_number": page_number,
                            "native_text_extracted": True,
                            "visual_analysis_status": "not_needed",
                            "caption_generated": False,
                        },
                    )
                )
                continue

            ocr_record = self._try_ocr_pdf_page(path, page_number)

            if ocr_record:
                records.append(ocr_record)
                continue

            if (
                self.vision_client.is_enabled()
                and self.vision_config.caption_pdf_pages
                and self.caption_budget.can_caption_pdf_page()
            ):
                self.caption_budget.use_pdf_page_caption()
                records.append(self._caption_pdf_page(path, page_number))
            else:
                records.append(self._catalog_image_only_pdf_page(path, page_number))

        return LoadedFile(
            source_path=path,
            records=records,
            document_metadata={
                "file_name": path.name,
                "file_suffix": path.suffix.lower(),
                "loader": "DocumentFileLoader",
                "page_count": len(reader.pages),
                "records_created": len(records),
                "vision_enabled": self.vision_client.is_enabled(),
            },
        )

    def _catalog_image_only_pdf_page(self, path: Path, page_number: int) -> LoadedRecord:
        """
        Create a metadata-only record for an image-only PDF page.
        """

        text = (
            f"PDF File: {path.name}\n"
            f"Page Number: {page_number}\n"
            "Native Text Extracted: False\n"
            "Visual Analysis Status: Not enabled\n"
            "Caption Generated: False\n"
            "Note: This PDF page appears to have no extractable embedded text. "
            "It may be scanned or image-based. A local vision model is required "
            "to visually interpret this page."
        )

        return LoadedRecord(
            text=text,
            record_type="pdf_page_image_metadata",
            title=f"{path.name} page {page_number}",
            metadata={
                "file_name": path.name,
                "file_suffix": path.suffix.lower(),
                "page_number": page_number,
                "native_text_extracted": False,
                "visual_analysis_status": "not_enabled",
                "caption_generated": False,
                "vision_provider": self.vision_config.provider,
                "vision_model": self.vision_config.model,
            },
        )

    def _try_ocr_pdf_page(self, path: Path, page_number: int) -> LoadedRecord | None:
        """
        Try OCR on an image-only PDF page before vision captioning.
        """

        if not (
            self.ocr_client.is_enabled()
            and self.ocr_config.ocr_pdf_pages
            and self.ocr_budget.can_ocr_pdf_page()
        ):
            return None

        self.ocr_budget.use_pdf_page_ocr()

        try:
            image_bytes = self._render_pdf_page_to_png_bytes(path, page_number)
            result = self.ocr_client.ocr_image_bytes(
                image_bytes=image_bytes,
                image_label=f"{path.name} page {page_number}",
            )

            if result.status != "ok":
                return None

            text = (
                f"PDF File: {path.name}\n"
                f"Page Number: {page_number}\n"
                "Native Text Extracted: False\n"
                "OCR Status: Extracted\n"
                "OCR Text:\n"
                f"{result.text}"
            )

            return LoadedRecord(
                text=text,
                record_type="pdf_page_ocr_text",
                title=f"{path.name} page {page_number}",
                metadata={
                    "file_name": path.name,
                    "file_suffix": path.suffix.lower(),
                    "page_number": page_number,
                    "native_text_extracted": False,
                    "ocr_status": "extracted",
                    "ocr_generated": True,
                    "ocr_provider": self.ocr_config.provider,
                    "ocr_languages": self.ocr_config.languages,
                },
            )

        except Exception:
            return None

    def _caption_pdf_page(self, path: Path, page_number: int) -> LoadedRecord:
        """
        Render and caption an image-only PDF page using local vision.
        """

        try:
            image_bytes = self._render_pdf_page_to_png_bytes(path, page_number)

            prompt = (
                "Describe this PDF page for a retrieval system. "
                "Extract any visible text if readable. Also describe visible products, "
                "equipment, diagrams, tables, labels, and important visual details. "
                "Do not invent details."
            )

            result = self.vision_client.caption_image_bytes(
                image_bytes=image_bytes,
                image_label=f"{path.name} page {page_number}",
                prompt=prompt,
            )

            if result.status != "ok":
                return self._catalog_pdf_page_with_vision_error(path, page_number, result.error)

            text = (
                f"PDF File: {path.name}\n"
                f"Page Number: {page_number}\n"
                "Native Text Extracted: False\n"
                "Visual Analysis Status: Captioned\n"
                "Caption Generated: True\n"
                f"PDF Page Caption:\n{result.caption}"
            )

            return LoadedRecord(
                text=text,
                record_type="pdf_page_vision_caption",
                title=f"{path.name} page {page_number}",
                metadata={
                    "file_name": path.name,
                    "file_suffix": path.suffix.lower(),
                    "page_number": page_number,
                    "native_text_extracted": False,
                    "visual_analysis_status": "captioned",
                    "caption_generated": True,
                    "vision_provider": self.vision_config.provider,
                    "vision_model": self.vision_config.model,
                },
            )

        except Exception as exc:
            return self._catalog_pdf_page_with_vision_error(path, page_number, str(exc))

    def _catalog_pdf_page_with_vision_error(
        self,
        path: Path,
        page_number: int,
        error: str | None,
    ) -> LoadedRecord:
        """
        Fallback if PDF page captioning fails.
        """

        base = self._catalog_image_only_pdf_page(path, page_number)
        metadata = dict(base.metadata)
        metadata["visual_analysis_status"] = "caption_failed"
        metadata["caption_error"] = error

        return LoadedRecord(
            text=base.text + "\nVision Caption Attempt: Failed\n" + f"Vision Error: {error}",
            record_type="pdf_page_image_metadata",
            title=base.title,
            metadata=metadata,
        )

    @staticmethod
    def _render_pdf_page_to_png_bytes(path: Path, page_number: int) -> bytes:
        """
        Render a PDF page to PNG bytes using PyMuPDF.

        page_number is 1-based.
        """

        document = fitz.open(str(path))

        try:
            page = document.load_page(page_number - 1)
            pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            return pixmap.tobytes("png")
        finally:
            document.close()

    def _load_docx(self, path: Path) -> LoadedFile:
        """
        Load a DOCX file as section-style records.
        """

        document = Document(str(path))
        records: list[LoadedRecord] = []

        current_title = path.name
        current_lines: list[str] = []
        section_index = 0

        def flush_section() -> None:
            nonlocal section_index, current_lines, current_title

            text = "\n".join(line for line in current_lines if line.strip()).strip()

            if not text:
                current_lines = []
                return

            records.append(
                LoadedRecord(
                    text=text,
                    record_type="docx_section",
                    title=current_title,
                    metadata={
                        "file_name": path.name,
                        "file_suffix": path.suffix.lower(),
                        "section_index": section_index,
                    },
                )
            )

            section_index += 1
            current_lines = []

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()

            if not text:
                continue

            style_name = paragraph.style.name if paragraph.style else ""

            if style_name.lower().startswith("heading"):
                flush_section()
                current_title = text
                current_lines.append(text)
            else:
                current_lines.append(text)

        table_text = self._extract_docx_tables(document)

        if table_text:
            current_lines.append("\nTables:")
            current_lines.extend(table_text)

        flush_section()

        return LoadedFile(
            source_path=path,
            records=records,
            document_metadata={
                "file_name": path.name,
                "file_suffix": path.suffix.lower(),
                "loader": "DocumentFileLoader",
                "section_count": len(records),
                "table_count": len(document.tables),
            },
        )

    def _load_pptx(self, path: Path) -> LoadedFile:
        """
        Load a PPTX file as one record per slide.
        """

        presentation = Presentation(str(path))
        records: list[LoadedRecord] = []

        for slide_index, slide in enumerate(presentation.slides):
            slide_number = slide_index + 1
            lines: list[str] = []

            for shape in slide.shapes:
                lines.extend(self._extract_text_from_pptx_shape(shape))

            text = "\n".join(line for line in lines if line.strip()).strip()

            if not text:
                continue

            title = self._first_non_empty_line(text) or f"{path.name} slide {slide_number}"

            records.append(
                LoadedRecord(
                    text=text,
                    record_type="pptx_slide",
                    title=title,
                    metadata={
                        "file_name": path.name,
                        "file_suffix": path.suffix.lower(),
                        "slide_number": slide_number,
                    },
                )
            )

        return LoadedFile(
            source_path=path,
            records=records,
            document_metadata={
                "file_name": path.name,
                "file_suffix": path.suffix.lower(),
                "loader": "DocumentFileLoader",
                "slide_count": len(presentation.slides),
                "records_created": len(records),
            },
        )

    @staticmethod
    def _extract_docx_tables(document: Any) -> list[str]:
        table_lines: list[str] = []

        for table_index, table in enumerate(document.tables):
            table_lines.append(f"Table {table_index + 1}:")

            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                cells = [cell for cell in cells if cell]

                if cells:
                    table_lines.append(" | ".join(cells))

        return table_lines

    @staticmethod
    def _extract_text_from_pptx_shape(shape: Any) -> list[str]:
        lines: list[str] = []

        if hasattr(shape, "text") and shape.text:
            lines.append(str(shape.text).strip())

        if hasattr(shape, "has_table") and shape.has_table:
            for row in shape.table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                cells = [cell for cell in cells if cell]

                if cells:
                    lines.append(" | ".join(cells))

        return lines

    @staticmethod
    def _first_non_empty_line(text: str) -> str | None:
        for line in text.splitlines():
            cleaned = line.strip()

            if cleaned:
                return cleaned

        return None